from pptx import Presentation
from pptx.exc import PythonPptxError
from sys import argv
from glob import glob
from unidecode import unidecode
from lxml.etree import XMLSyntaxError

for eachfile in glob(pathname="**/*.pptx", recursive=True):
    # skip files which contains ~ or $
    if eachfile[0][0] == '~' or \
        eachfile[0][0] =='$' or \
        '/~' in eachfile or \
        '/$' in eachfile:
        continue

    print("*** "+eachfile)
    try:    
        prs = Presentation(eachfile)
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape,"text"):
                    if unidecode(argv[1]).lower() in unidecode(shape.text).lower():
                        print(shape.text)
    except (PythonPptxError, XMLSyntaxError):
        print("### file "+eachfile+" not processed (pptx error)")
        continue
